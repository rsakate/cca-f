#!/usr/bin/env python3
"""Generate CCA-F Module 4 PowerPoint presentation — Prompt Engineering & Structured Output."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
ACCENT_TEAL = RGBColor(0x1A, 0xBC, 0x9C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
MED_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
SECTION_BG = RGBColor(0xF8, 0xF9, 0xFA)
CODE_BG = RGBColor(0x2D, 0x2D, 0x2D)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helper functions ───

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=DARK_GRAY, font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, text_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color if text_color else color
        p.font.bold = is_bold
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
    return txBox

def add_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_down_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color, text="", font_size=14, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shape

def add_rounded_box_with_text(slide, left, top, width, height, fill_color, text, font_size=14, text_color=WHITE, border_color=None):
    shape = add_shape_box(slide, left, top, width, height, fill_color, border_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Calibri"
    return shape

def add_connector_line(slide, x1, y1, x2, y2, color=MED_GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def section_header_slide(title, subtitle, color=ACCENT_BLUE, section_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, color)
    if section_num:
        add_circle(slide, Inches(5.9), Inches(1.5), Inches(1.5), WHITE, section_num, 36, color)
    add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1.5), title, 44, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2), Inches(4.5), Inches(9), Inches(1), subtitle, 22, RGBColor(0xDD, 0xDD, 0xDD), False, PP_ALIGN.CENTER)
    return slide

def content_slide(title, section_label=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.7), title, 30, DARK_GRAY, True)
    if section_label:
        add_textbox(slide, Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5), section_label, 14, MED_GRAY, False, PP_ALIGN.RIGHT)
    return slide


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2), "CCA-F Module 4", 52, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1), "Prompt Engineering & Structured Output", 36, ACCENT_BLUE, False, PP_ALIGN.CENTER)
add_shape_box(slide, Inches(4.5), Inches(4.0), Inches(4.3), Inches(0.04), ACCENT_BLUE)
add_textbox(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.7), "Claude Certified Architect -- Foundations", 22, MED_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5), "Quick Reference Study Deck", 18, MED_GRAY, False, PP_ALIGN.CENTER)

topics = ["S1: Explicit\nCriteria", "S2: Few-Shot\nPrompting", "S3: tool_use &\nJSON Schema", "S4: Validation\n& Retry Loops", "S5: Batch\nProcessing", "S6: Multi-Pass\nArchitectures"]
colors = [ACCENT_RED, ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_TEAL]
for i, (topic, clr) in enumerate(zip(topics, colors)):
    x = Inches(0.4 + i * 2.15)
    add_rounded_box_with_text(slide, x, Inches(6.0), Inches(1.95), Inches(0.9), clr, topic, 11, WHITE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Module Overview / Roadmap
# ═══════════════════════════════════════════════════════════════════
slide = content_slide("Module 4 -- Learning Roadmap")

steps = ["Self-Paced\nLearning", "Labs", "Instructor\nLed Session", "Assessment", "Completion"]
step_colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_PURPLE, ACCENT_RED]
for i, (step, clr) in enumerate(zip(steps, step_colors)):
    x = Inches(0.8 + i * 2.5)
    add_rounded_box_with_text(slide, x, Inches(1.5), Inches(2.0), Inches(0.8), clr, step, 13, WHITE)
    if i < 4:
        add_arrow(slide, Inches(2.8 + i * 2.5), Inches(1.7), Inches(0.5), Inches(0.35), MED_GRAY)

add_textbox(slide, Inches(0.8), Inches(2.8), Inches(5), Inches(0.5), "Builds on Skilljar Course:", 20, DARK_GRAY, True)
courses = [
    ("1", "Building with the Claude API", "Prompt techniques, tool_use, structured data, workflows"),
]
for i, (num, name, desc) in enumerate(courses):
    y = Inches(3.4 + i * 0.7)
    add_circle(slide, Inches(0.9), y, Inches(0.45), ACCENT_BLUE, num, 14, WHITE)
    add_textbox(slide, Inches(1.6), y, Inches(5), Inches(0.35), name, 15, DARK_GRAY, True)
    add_textbox(slide, Inches(1.6), y + Inches(0.3), Inches(5), Inches(0.3), desc, 11, MED_GRAY)

add_textbox(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.5), "3 Labs (2 mandatory, 1 optional):", 18, DARK_GRAY, True)
labs = [
    ("4.1", "Precision Prompting", "Explicit criteria & few-shot consistency", ACCENT_RED),
    ("4.2", "Enforcing Structure", "tool_use schemas + validation & retry", ACCENT_GREEN),
    ("4.3", "Scaling Output (Optional)", "Batch processing & multi-pass review", ACCENT_PURPLE),
]
for i, (num, name, desc, clr) in enumerate(labs):
    y = Inches(5.0 + i * 0.7)
    add_rounded_box_with_text(slide, Inches(0.9), y, Inches(0.6), Inches(0.4), clr, num, 10, WHITE)
    add_textbox(slide, Inches(1.7), y, Inches(2.5), Inches(0.35), name, 14, DARK_GRAY, True)
    add_textbox(slide, Inches(4.3), y, Inches(3), Inches(0.35), desc, 11, MED_GRAY)

add_textbox(slide, Inches(7.5), Inches(2.8), Inches(5), Inches(0.5), "6 Key Sections:", 20, DARK_GRAY, True)
sections = [
    ("S1", "Designing Prompts with Explicit Criteria", "45 min"),
    ("S2", "Few-Shot Prompting", "30 min"),
    ("S3", "Structured Output (tool_use + JSON)", "60 min"),
    ("S4", "Validation, Retry & Feedback Loops", "60 min"),
    ("S5", "Batch Processing with Claude API", "30 min"),
    ("S6", "Multi-Instance & Multi-Pass", "45 min"),
]
for i, (sec, name, time) in enumerate(sections):
    y = Inches(3.4 + i * 0.6)
    add_rounded_box_with_text(slide, Inches(7.5), y, Inches(0.6), Inches(0.4), colors[i], sec, 10, WHITE)
    add_textbox(slide, Inches(8.3), y, Inches(3), Inches(0.4), name, 14, DARK_GRAY, True)
    add_textbox(slide, Inches(11.5), y, Inches(1.2), Inches(0.4), time, 12, MED_GRAY, False, PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# S1: DESIGNING PROMPTS WITH EXPLICIT CRITERIA
# ═══════════════════════════════════════════════════════════════════
section_header_slide("S1: Designing Prompts with\nExplicit Criteria", "Clear, categorical rules and severity tiers to cut false positives", ACCENT_RED, "S1")


# --- S1 Slide: What Are Explicit Criteria? ---
slide = content_slide("What Are Explicit Criteria?", "S1")

add_shape_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.5), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
intro_lines = [
    ("Explicit criteria = Concrete, testable rules that tell Claude EXACTLY what to do.", True, DARK_GRAY),
    ("Instead of 'find bugs', say 'flag as critical if: SQL injection, unsanitized user input, or hardcoded secrets'.", False, DARK_GRAY),
    ("Think of it like giving a judge a rulebook instead of saying 'be fair'.", False, MED_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.3), intro_lines, 16)

add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5), "The Three Pillars of Explicit Prompting:", 20, DARK_GRAY, True)

pillars = [
    ("CATEGORICAL\nRULES", "Define exact categories with\nclear boundaries\n\nNot: 'find issues'\nBut: 'classify as: bug,\nperformance, security, style'",
     "Like a SORTING HAT\nin Harry Potter -- each\nitem goes to exactly\none house", ACCENT_RED),
    ("SEVERITY\nTIERS", "Define what each level means\nwith concrete examples\n\nCritical = data loss, security\nHigh = feature broken\nMedium = degraded UX",
     "Like a HOSPITAL\nTRIAGE system --\nred, yellow, green\ntags on patients", ACCENT_BLUE),
    ("DECISION\nBOUNDARIES", "Define the edge cases explicitly\nso Claude doesn't guess\n\n'If unsure between High and\nMedium, default to High'",
     "Like SPEED LIMITS --\nclear numbers, not\n'drive reasonably'", ACCENT_GREEN),
]

for i, (name, desc, analogy, clr) in enumerate(pillars):
    x = Inches(0.5 + i * 4.2)
    add_rounded_box_with_text(slide, x, Inches(3.6), Inches(3.7), Inches(0.7), clr, name, 16, WHITE)

    add_shape_box(slide, x, Inches(4.4), Inches(3.7), Inches(1.8), LIGHT_GRAY, clr, Pt(1))
    add_textbox(slide, x + Inches(0.2), Inches(4.5), Inches(3.3), Inches(1.6), desc, 12, DARK_GRAY)

    add_shape_box(slide, x, Inches(6.4), Inches(3.7), Inches(0.9), RGBColor(0xFD, 0xF2, 0xE9), clr, Pt(1))
    add_textbox(slide, x + Inches(0.2), Inches(6.42), Inches(3.3), Inches(0.85), analogy, 11, clr, True, PP_ALIGN.CENTER)


# --- S1 Slide: Vague vs Explicit Prompts ---
slide = content_slide("Vague vs Explicit Prompts -- Side by Side", "S1")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "The #1 prompt engineering mistake: being vague. Compare these two approaches:", 17, DARK_GRAY, True)

# BAD: Vague
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(3.0), RGBColor(0xFA, 0xEB, 0xEB), ACCENT_RED, Pt(2))
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(0.4), "BAD: Vague Prompt", 20, ACCENT_RED, True)
bad_lines = [
    ('"Review this code and find any', False, WHITE),
    (' issues that might be important."', False, WHITE),
    ("", False, None),
    ("Problems:", True, ACCENT_RED),
    ("- What counts as 'important'?", False, MED_GRAY),
    ("- No categories defined", False, MED_GRAY),
    ("- No severity levels", False, MED_GRAY),
    ("- Claude guesses -> inconsistent", False, MED_GRAY),
    ("- High false positive rate", False, MED_GRAY),
    ("- Different results each run!", False, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.3), Inches(2.4), bad_lines, 13, DARK_GRAY, "Courier New")

# GOOD: Explicit
add_shape_box(slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(3.0), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.4), "GOOD: Explicit Prompt", 20, ACCENT_GREEN, True)
good_lines = [
    ('"Classify each issue as:', False, WHITE),
    (' SECURITY | PERFORMANCE | STYLE', False, ACCENT_GREEN),
    (' Severity: CRITICAL if data loss', False, ACCENT_GREEN),
    (' or injection. HIGH if feature', False, ACCENT_GREEN),
    (' breaks. MEDIUM if UX degrades."', False, ACCENT_GREEN),
    ("", False, None),
    ("Benefits:", True, ACCENT_GREEN),
    ("- Clear categories = consistent", False, MED_GRAY),
    ("- Defined severity = actionable", False, MED_GRAY),
    ("- Reproducible across runs!", False, ACCENT_GREEN),
]
add_multiline_textbox(slide, Inches(7.1), Inches(2.4), Inches(5.3), Inches(2.4), good_lines, 13, DARK_GRAY, "Courier New")

# Real-life analogy
add_shape_box(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.1), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.4), "Real-Life Analogy: Grading Student Essays", 16, ACCENT_ORANGE, True)
analogy_lines = [
    ("Vague: 'Grade these essays fairly.' -- Every teacher grades differently. No consistency!", False, ACCENT_RED),
    ("Explicit: 'Use this rubric: Thesis (20pts), Evidence (30pts), Grammar (20pts), Structure (30pts)'.", False, ACCENT_GREEN),
    ("With a rubric, ANY teacher produces consistent grades. Explicit prompts are your rubric for Claude!", False, ACCENT_ORANGE),
    ("", False, None),
    ("Reducing false positives = iterating on your rubric. Run it, check results, tighten the criteria.", True, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.5), analogy_lines, 14)


# --- S1 Slide: Severity Tiers in Practice ---
slide = content_slide("Designing Severity Tiers -- A Practical Example", "S1")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "Severity tiers turn subjective judgment into objective, repeatable classification:", 17, DARK_GRAY, True)

# Code example
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(6.5), Inches(4.5), CODE_BG)
code_lines = [
    ("# Prompt with explicit severity tiers", False, MED_GRAY),
    ("", False, None),
    ('system_prompt = """', False, WHITE),
    ("Classify each code issue found.", False, WHITE),
    ("", False, None),
    ("SEVERITY DEFINITIONS:", True, ACCENT_GREEN),
    ("", False, None),
    ("CRITICAL:", True, ACCENT_RED),
    ("  - SQL injection or XSS", False, RGBColor(0xF0, 0xC6, 0x74)),
    ("  - Hardcoded secrets/passwords", False, RGBColor(0xF0, 0xC6, 0x74)),
    ("  - Data loss or corruption risk", False, RGBColor(0xF0, 0xC6, 0x74)),
    ("", False, None),
    ("HIGH:", True, ACCENT_ORANGE),
    ("  - Feature completely broken", False, RGBColor(0xF0, 0xC6, 0x74)),
    ("  - Authentication bypass", False, RGBColor(0xF0, 0xC6, 0x74)),
    ("", False, None),
    ("MEDIUM:", True, ACCENT_BLUE),
    ("  - Degraded user experience", False, RGBColor(0xF0, 0xC6, 0x74)),
    ("  - Missing error handling", False, RGBColor(0xF0, 0xC6, 0x74)),
    ("", False, None),
    ("LOW:", True, ACCENT_GREEN),
    ("  - Code style / naming issues", False, RGBColor(0xF0, 0xC6, 0x74)),
    ('"""', False, WHITE),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.0), Inches(6.1), Inches(4.3), code_lines, 12, WHITE, "Courier New")

# Key points
add_textbox(slide, Inches(7.5), Inches(1.9), Inches(5.5), Inches(0.5), "Why This Works:", 20, DARK_GRAY, True)

points = [
    ("Concrete Examples", "Each tier has specific, testable\ncriteria -- no guessing needed", ACCENT_RED),
    ("Clear Boundaries", "If unsure between tiers, the\nexamples guide the decision", ACCENT_BLUE),
    ("Iterate to Improve", "Run the prompt, check false\npositives, tighten criteria", ACCENT_GREEN),
    ("Consistent Output", "Same input = same severity\nacross multiple runs", ACCENT_ORANGE),
]
for i, (title, desc, clr) in enumerate(points):
    y = Inches(2.5 + i * 1.15)
    add_rounded_box_with_text(slide, Inches(7.5), y, Inches(2.5), Inches(0.4), clr, title, 13, WHITE)
    add_textbox(slide, Inches(7.5), y + Inches(0.5), Inches(5.5), Inches(0.6), desc, 13, DARK_GRAY)

# Exam tip
add_shape_box(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6), RGBColor(0xFA, 0xEB, 0xEB), ACCENT_RED, Pt(2))
add_textbox(slide, Inches(0.9), Inches(6.72), Inches(11.5), Inches(0.5),
    "EXAM TIP: Explicit criteria = categorical rules + severity tiers + concrete examples. This is how you cut false positives.", 15, ACCENT_RED, True)


# ═══════════════════════════════════════════════════════════════════
# S2: FEW-SHOT PROMPTING FOR CONSISTENCY & GENERALIZATION
# ═══════════════════════════════════════════════════════════════════
section_header_slide("S2: Few-Shot Prompting for\nConsistency & Generalization", "Use 2-4 labeled examples to lock in format, tone, and edge-case handling", ACCENT_BLUE, "S2")


# --- S2 Slide: What is Few-Shot Prompting? ---
slide = content_slide("What is Few-Shot Prompting?", "S2")

add_shape_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.5), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
intro_lines = [
    ("Few-shot prompting = Giving Claude 2-4 worked examples BEFORE asking it to do the task.", True, DARK_GRAY),
    ("The examples show the expected input, output format, and decision logic for ambiguous cases.", False, DARK_GRAY),
    ("Like training a new employee by showing them 3 completed reports before asking them to write one.", False, MED_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.3), intro_lines, 16)

# Shot types comparison
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5), "The Prompting Spectrum:", 20, DARK_GRAY, True)

shot_types = [
    ("ZERO-SHOT", "No examples at all.\nJust instructions.",
     "Works for simple tasks\nwhere the format is obvious.\n\n'Translate this to French'",
     "Like saying 'cook\ndinner' to a chef", ACCENT_RED),
    ("FEW-SHOT\n(2-4 examples)", "Provide labeled examples\nthat demonstrate format.",
     "Best for: complex extraction,\nambiguous classification, tasks\nwhere format matters.\n\n'Here are 3 examples...'",
     "Like showing a chef\n3 photos of the dish\nyou want", ACCENT_GREEN),
    ("MANY-SHOT\n(10+ examples)", "Large set of examples.\nHigher cost, diminishing returns.",
     "Rarely needed. Use only if\nfew-shot isn't consistent\nenough. Eats tokens fast.\n\nUsually 2-4 is plenty!",
     "Like giving a chef\nan entire cookbook\n(overkill!)", ACCENT_ORANGE),
]

for i, (name, subtitle, desc, analogy, clr) in enumerate(shot_types):
    x = Inches(0.5 + i * 4.2)
    add_rounded_box_with_text(slide, x, Inches(3.5), Inches(3.7), Inches(0.7), clr, name, 16, WHITE)
    add_textbox(slide, x + Inches(0.1), Inches(4.3), Inches(3.5), Inches(0.4), subtitle, 12, DARK_GRAY, True, PP_ALIGN.CENTER)

    add_shape_box(slide, x, Inches(4.8), Inches(3.7), Inches(1.6), LIGHT_GRAY, clr, Pt(1))
    add_textbox(slide, x + Inches(0.2), Inches(4.9), Inches(3.3), Inches(1.4), desc, 12, DARK_GRAY)

    add_shape_box(slide, x, Inches(6.6), Inches(3.7), Inches(0.7), RGBColor(0xFD, 0xF2, 0xE9), clr, Pt(1))
    add_textbox(slide, x + Inches(0.2), Inches(6.62), Inches(3.3), Inches(0.65), analogy, 11, clr, True, PP_ALIGN.CENTER)


# --- S2 Slide: Few-Shot Example in Code ---
slide = content_slide("Few-Shot Prompting -- Code Example", "S2")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "Show Claude the pattern with examples, then give the real task. It matches the format automatically:", 17, DARK_GRAY, True)

# Code example
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(7.0), Inches(5.0), CODE_BG)
code_lines = [
    ("# Few-shot: sentiment classification", False, MED_GRAY),
    ("", False, None),
    ("messages = [", False, WHITE),
    ('  # Example 1 - Clear positive', False, MED_GRAY),
    ('  {"role":"user","content":', False, WHITE),
    ('   "Review: Great product!"},', False, ACCENT_BLUE),
    ('  {"role":"assistant","content":', False, WHITE),
    ('   "sentiment:positive|conf:0.95"},', False, ACCENT_GREEN),
    ("", False, None),
    ('  # Example 2 - Tricky sarcasm', False, MED_GRAY),
    ('  {"role":"user","content":', False, WHITE),
    ('   "Review: Oh sure, it works GREAT', False, ACCENT_BLUE),
    ('    if you enjoy waiting 10 min"},', False, ACCENT_BLUE),
    ('  {"role":"assistant","content":', False, WHITE),
    ('   "sentiment:negative|conf:0.85"},', False, ACCENT_RED),
    ("", False, None),
    ('  # Example 3 - Genuinely mixed', False, MED_GRAY),
    ('  {"role":"user","content":', False, WHITE),
    ('   "Review: Fast shipping, but item', False, ACCENT_BLUE),
    ('    was damaged on arrival"},', False, ACCENT_BLUE),
    ('  {"role":"assistant","content":', False, WHITE),
    ('   "sentiment:mixed|conf:0.80"},', False, ACCENT_ORANGE),
    ("", False, None),
    ('  # NOW the real task', False, MED_GRAY),
    ('  {"role":"user","content":', False, WHITE),
    ('   "Review: Decent for the price"},', False, ACCENT_BLUE),
    ("]", False, WHITE),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.0), Inches(6.6), Inches(4.8), code_lines, 11, WHITE, "Courier New")

# Key insights
add_textbox(slide, Inches(8.0), Inches(1.9), Inches(5.0), Inches(0.5), "Why These 3 Examples Work:", 18, DARK_GRAY, True)

insights = [
    ("Cover edge cases", "Example 2 handles sarcasm --\nthe hardest case. Claude now\nknows HOW to handle it.", ACCENT_RED),
    ("Lock in format", "The pipe-delimited format is\nconsistent across all examples.\nClaude matches it exactly.", ACCENT_BLUE),
    ("Show confidence", "Each example includes a\nconfidence score. Claude learns\nto self-assess certainty.", ACCENT_GREEN),
    ("Generalize", "Claude can now classify unseen\nreviews by PATTERN MATCHING\nthe logic from examples.", ACCENT_ORANGE),
]
for i, (title, desc, clr) in enumerate(insights):
    y = Inches(2.5 + i * 1.2)
    add_rounded_box_with_text(slide, Inches(8.0), y, Inches(2.3), Inches(0.4), clr, title, 12, WHITE)
    add_textbox(slide, Inches(8.0), y + Inches(0.5), Inches(5.0), Inches(0.7), desc, 12, DARK_GRAY)

# Bottom tip
add_shape_box(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4), RGBColor(0xFA, 0xEB, 0xEB), ACCENT_RED, Pt(1))
add_textbox(slide, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
    "EXAM TIP: Few-shot examples fix tone + format. Cover AMBIGUOUS cases (sarcasm, edge cases) -- easy ones don't need examples!", 13, ACCENT_RED, True)


# --- S2 Slide: Best Practices ---
slide = content_slide("Few-Shot Best Practices & Generalization", "S2")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "How many examples? Which ones? And how does Claude generalize beyond them?", 17, DARK_GRAY, True)

# How many
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(3.7), Inches(3.0), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(3.3), Inches(0.4), "How Many Examples?", 18, ACCENT_BLUE, True)
how_many = [
    ("2 examples: Minimum. Shows", False, DARK_GRAY),
    ("   the format but limited coverage.", False, DARK_GRAY),
    ("", False, None),
    ("3-4 examples: Sweet spot.", True, ACCENT_GREEN),
    ("   Enough to show pattern +", False, DARK_GRAY),
    ("   edge cases without wasting", False, DARK_GRAY),
    ("   tokens.", False, DARK_GRAY),
    ("", False, None),
    ("5+: Diminishing returns. Only", False, DARK_GRAY),
    ("   if task is very complex.", False, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.4), Inches(3.3), Inches(2.4), how_many, 12)

# Which cases
add_shape_box(slide, Inches(4.8), Inches(1.9), Inches(3.7), Inches(3.0), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(5.0), Inches(2.0), Inches(3.3), Inches(0.4), "Which Cases to Include?", 18, ACCENT_GREEN, True)
which_cases = [
    ("DO include:", True, ACCENT_GREEN),
    ("- Ambiguous / edge cases", False, DARK_GRAY),
    ("- Cases where format matters", False, DARK_GRAY),
    ("- Tricky decisions (sarcasm)", False, DARK_GRAY),
    ("", False, None),
    ("DON'T include:", True, ACCENT_RED),
    ("- Obvious, simple cases", False, DARK_GRAY),
    ("- Redundant similar examples", False, DARK_GRAY),
    ("- Cases with no decision point", False, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(5.0), Inches(2.4), Inches(3.3), Inches(2.4), which_cases, 12)

# Generalization
add_shape_box(slide, Inches(9.0), Inches(1.9), Inches(3.7), Inches(3.0), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(9.2), Inches(2.0), Inches(3.3), Inches(0.4), "Generalization", 18, ACCENT_ORANGE, True)
generalize = [
    ("Claude doesn't memorize examples.", True, DARK_GRAY),
    ("It learns the PATTERN:", False, DARK_GRAY),
    ("", False, None),
    ("- Input format -> Output format", False, ACCENT_BLUE),
    ("- Decision logic for edge cases", False, ACCENT_GREEN),
    ("- Confidence calibration", False, ACCENT_ORANGE),
    ("", False, None),
    ("Test with UNSEEN inputs to", True, ACCENT_RED),
    ("verify it generalizes well!", False, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(9.2), Inches(2.4), Inches(3.3), Inches(2.4), generalize, 12)

# Analogy
add_shape_box(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.1), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.4), "Real-Life Analogy: Training a New Barista", 16, ACCENT_ORANGE, True)
barista_lines = [
    ("You're training a new barista at Starbucks. You don't make them memorize 10,000 orders.", False, DARK_GRAY),
    ("You show them 3-4 examples: a simple latte, a complex custom order, a tricky modification, and a refund request.", False, DARK_GRAY),
    ("From those 4 examples, they GENERALIZE to handle any order that comes in -- even ones you didn't show them!", False, ACCENT_GREEN),
    ("", False, None),
    ("Few-shot works the same: show Claude the PATTERN with diverse examples, and it handles the rest.", True, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.5), barista_lines, 14)


# ═══════════════════════════════════════════════════════════════════
# S3: STRUCTURED OUTPUT USING tool_use & JSON SCHEMA
# ═══════════════════════════════════════════════════════════════════
section_header_slide("S3: Structured Output using\ntool_use & JSON Schema", "Force Claude to return valid, parseable JSON every time -- not free-text", ACCENT_GREEN, "S3")


# --- S3 Slide: What is tool_use for Structured Output? ---
slide = content_slide("What is tool_use for Structured Output?", "S3")

add_shape_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.5), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
intro_lines = [
    ("tool_use = A Claude API feature that forces the model to return structured data matching a JSON schema.", True, DARK_GRAY),
    ("Instead of Claude writing free-text, it fills in a structured 'form' with typed fields you define.", False, DARK_GRAY),
    ("Think of it as giving Claude a form to fill out instead of a blank piece of paper.", False, MED_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.3), intro_lines, 16)

# How it works
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5), "How tool_use Forces Structured Output:", 20, DARK_GRAY, True)

# Flow diagram
flow_steps = [
    ("You Define\na Tool Schema", "JSON Schema with\ntyped fields", ACCENT_BLUE),
    ("Set tool_choice\nto 'any'", "Forces Claude to\nuse the tool", ACCENT_GREEN),
    ("Claude Returns\nStructured JSON", "Matches your\nschema exactly", ACCENT_ORANGE),
    ("Parse & Use\nthe Data", "Guaranteed valid\nJSON to process", ACCENT_PURPLE),
]
for i, (name, desc, clr) in enumerate(flow_steps):
    x = Inches(0.5 + i * 3.2)
    add_rounded_box_with_text(slide, x, Inches(3.6), Inches(2.8), Inches(0.7), clr, name, 13, WHITE)
    add_textbox(slide, x + Inches(0.1), Inches(4.4), Inches(2.6), Inches(0.5), desc, 12, DARK_GRAY, False, PP_ALIGN.CENTER)
    if i < 3:
        add_arrow(slide, Inches(3.3 + i * 3.2), Inches(3.8), Inches(0.4), Inches(0.3), MED_GRAY)

# Free-text vs structured comparison
add_shape_box(slide, Inches(0.6), Inches(5.2), Inches(5.8), Inches(2.1), RGBColor(0xFA, 0xEB, 0xEB), ACCENT_RED, Pt(2))
add_textbox(slide, Inches(0.8), Inches(5.3), Inches(5.3), Inches(0.4), "Without tool_use (free-text):", 16, ACCENT_RED, True)
free_lines = [
    ('"The bug is critical, it\'s a SQL', False, MED_GRAY),
    (' injection in the login form and', False, MED_GRAY),
    (' the severity is high..."', False, MED_GRAY),
    ("", False, None),
    ("Parsing this is fragile and error-prone!", True, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(0.8), Inches(5.7), Inches(5.3), Inches(1.5), free_lines, 12, DARK_GRAY, "Courier New")

add_shape_box(slide, Inches(6.9), Inches(5.2), Inches(5.8), Inches(2.1), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(7.1), Inches(5.3), Inches(5.3), Inches(0.4), "With tool_use (structured):", 16, ACCENT_GREEN, True)
struct_lines = [
    ('{"category": "security",', False, ACCENT_GREEN),
    (' "severity": "critical",', False, ACCENT_GREEN),
    (' "location": "login.py:42",', False, ACCENT_GREEN),
    (' "description": "SQL injection"}', False, ACCENT_GREEN),
    ("", False, None),
    ("Clean JSON. Parse with json.loads()!", True, ACCENT_GREEN),
]
add_multiline_textbox(slide, Inches(7.1), Inches(5.7), Inches(5.3), Inches(1.5), struct_lines, 12, DARK_GRAY, "Courier New")


# --- S3 Slide: tool_choice Options ---
slide = content_slide("tool_choice Options -- Controlling When Tools Run", "S3")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "tool_choice controls WHETHER and WHICH tool Claude uses. Critical for forcing structured output:", 17, DARK_GRAY, True)

# Four options
options = [
    ("auto", "Claude DECIDES whether\nto use a tool or not",
     "Default behavior. Claude may\nreturn free text OR use a tool.\nUse when both are acceptable.",
     "FLEXIBLE", ACCENT_BLUE),
    ("any", "Claude MUST use\none of the tools",
     "Forces structured output!\nClaude picks which tool but\nMUST return tool output.",
     "FORCE STRUCTURE", ACCENT_GREEN),
    ('{"type":"tool",\n"name":"X"}', "Claude MUST use\nthe SPECIFIC tool X",
     "Forces a specific tool.\nUse when you have multiple\ntools and want a specific one.",
     "FORCE SPECIFIC", ACCENT_ORANGE),
    ("none", "Claude CANNOT\nuse any tools",
     "Disabled. Claude only returns\nfree text even if tools are\ndefined. Rarely used.",
     "DISABLED", MED_GRAY),
]

for i, (name, subtitle, desc, badge, clr) in enumerate(options):
    x = Inches(0.3 + i * 3.25)
    add_rounded_box_with_text(slide, x, Inches(1.8), Inches(3.0), Inches(0.6), clr, name, 14, WHITE)
    add_textbox(slide, x + Inches(0.1), Inches(2.5), Inches(2.8), Inches(0.5), subtitle, 11, DARK_GRAY, True, PP_ALIGN.CENTER)

    add_shape_box(slide, x, Inches(3.1), Inches(3.0), Inches(1.5), LIGHT_GRAY, clr, Pt(1))
    add_textbox(slide, x + Inches(0.15), Inches(3.2), Inches(2.7), Inches(1.3), desc, 11, DARK_GRAY)

    add_rounded_box_with_text(slide, x + Inches(0.3), Inches(4.8), Inches(2.4), Inches(0.4), clr, badge, 11, WHITE)

# Key insight box
add_shape_box(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.8), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(5.6), Inches(11), Inches(0.4), "Real-Life Analogy: Restaurant Ordering", 16, ACCENT_BLUE, True)
restaurant_lines = [
    ("auto = 'Order whatever you want' (might get a meal, might just get a conversation)", False, ACCENT_BLUE),
    ("any = 'You MUST order from the menu' (guaranteed to get a dish, you pick which one)", False, ACCENT_GREEN),
    ('tool:specific = \'You MUST order the Caesar salad\' (guaranteed specific dish)', False, ACCENT_ORANGE),
    ("none = 'No ordering allowed, just chat' (no food at all)", False, MED_GRAY),
    ("", False, None),
    ("For structured output, use tool_choice: 'any' or a specific tool name. NEVER 'auto' if you need guaranteed JSON!", True, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.2), restaurant_lines, 13)


# --- S3 Slide: Schema Design Best Practices ---
slide = content_slide("JSON Schema Design -- nullable, enum, and Detail Fields", "S3")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "A well-designed schema catches bad data at the structure level. Key patterns:", 17, DARK_GRAY, True)

# Code example
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(6.5), Inches(5.0), CODE_BG)
schema_code = [
    ("# Tool schema with best practices", False, MED_GRAY),
    ("", False, None),
    ("tool = {", False, WHITE),
    ('  "name": "classify_issue",', False, WHITE),
    ('  "description": "Classify a code issue",', False, WHITE),
    ('  "input_schema": {', False, WHITE),
    ('    "type": "object",', False, WHITE),
    ('    "properties": {', False, WHITE),
    ("", False, None),
    ('      "category": {', False, WHITE),
    ('        "type": "string",', False, WHITE),
    ('        "enum": ["security",', True, ACCENT_GREEN),
    ('          "performance","style"]', True, ACCENT_GREEN),
    ("      },", False, WHITE),
    ("", False, None),
    ('      "severity": {', False, WHITE),
    ('        "type": "string",', False, WHITE),
    ('        "enum": ["critical",', True, ACCENT_GREEN),
    ('          "high","medium","low"]', True, ACCENT_GREEN),
    ("      },", False, WHITE),
    ("", False, None),
    ('      "fix_suggestion": {', False, WHITE),
    ('        "type": ["string","null"]', True, ACCENT_ORANGE),
    ("      },", False, WHITE),
    ("", False, None),
    ('      "confidence": {', False, WHITE),
    ('        "type": "number"', False, WHITE),
    ("      }", False, WHITE),
    ("    },", False, WHITE),
    ('    "required": ["category",', False, ACCENT_RED),
    ('      "severity","confidence"]', False, ACCENT_RED),
    ("  }", False, WHITE),
    ("}", False, WHITE),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.0), Inches(6.1), Inches(4.8), schema_code, 11, WHITE, "Courier New")

# Patterns explained
add_textbox(slide, Inches(7.5), Inches(1.9), Inches(5.5), Inches(0.5), "Key Schema Patterns:", 20, DARK_GRAY, True)

patterns = [
    ("enum fields", 'Restrict to valid values ONLY.\n"enum": ["security","performance"]\nprevents Claude inventing new\ncategories like "bug" or "issue".', ACCENT_GREEN),
    ("nullable fields", '"type": ["string", "null"]\nAllows the field to be null when\ndata is genuinely missing.\nBetter than empty string ""!', ACCENT_ORANGE),
    ("required fields", '"required": ["category"]\nForces Claude to always include\nthese fields. Prevents partial\nor incomplete responses.', ACCENT_RED),
]
for i, (title, desc, clr) in enumerate(patterns):
    y = Inches(2.5 + i * 1.5)
    add_rounded_box_with_text(slide, Inches(7.5), y, Inches(2.3), Inches(0.4), clr, title, 13, WHITE)
    add_textbox(slide, Inches(7.5), y + Inches(0.5), Inches(5.5), Inches(1.0), desc, 12, DARK_GRAY)

# Important callout
add_shape_box(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4), RGBColor(0xFA, 0xEB, 0xEB), ACCENT_RED, Pt(1))
add_textbox(slide, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
    "EXAM TIP: Schema catches SYNTAX errors (wrong type, missing field). It CANNOT catch LOGIC errors (score=11 when max is 10)!", 13, ACCENT_RED, True)


# ═══════════════════════════════════════════════════════════════════
# S4: VALIDATION, RETRY & FEEDBACK LOOPS
# ═══════════════════════════════════════════════════════════════════
section_header_slide("S4: Validation, Retry &\nFeedback Loops", "Schema validation catches syntax, semantic checks catch logic -- combine both for reliability", ACCENT_ORANGE, "S4")


# --- S4 Slide: Schema vs Semantic Validation ---
slide = content_slide("Two Layers of Validation: Schema + Semantic", "S4")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "A JSON schema alone isn't enough. You need TWO layers of validation:", 17, DARK_GRAY, True)

# Schema validation
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(3.0), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(0.4), "Layer 1: Schema Validation (Syntax)", 18, ACCENT_BLUE, True)
schema_lines = [
    ("Checks the STRUCTURE of the data:", True, DARK_GRAY),
    ("", False, None),
    ("- Is it valid JSON?", False, MED_GRAY),
    ("- Are required fields present?", False, MED_GRAY),
    ("- Are types correct (string, number)?", False, MED_GRAY),
    ("- Do enum values match allowed set?", False, MED_GRAY),
    ("", False, None),
    ("Catches: missing fields, wrong types,", True, ACCENT_BLUE),
    ("invalid enum values, malformed JSON", False, ACCENT_BLUE),
    ("", False, None),
    ("Cannot catch: 'score: 99' when max=10", True, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.3), Inches(2.5), schema_lines, 13)

# Semantic validation
add_shape_box(slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(3.0), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.4), "Layer 2: Semantic Validation (Logic)", 18, ACCENT_GREEN, True)
semantic_lines = [
    ("Checks the MEANING of the data:", True, DARK_GRAY),
    ("", False, None),
    ("- Is score between 0-10?", False, MED_GRAY),
    ("- Does end_date come after start_date?", False, MED_GRAY),
    ("- Is the email format valid?", False, MED_GRAY),
    ("- Do cross-field constraints hold?", False, MED_GRAY),
    ("", False, None),
    ("Catches: out-of-range values, logical", True, ACCENT_GREEN),
    ("contradictions, business rule violations", False, ACCENT_GREEN),
    ("", False, None),
    ("Needs custom code -- schema can't do it!", True, ACCENT_ORANGE),
]
add_multiline_textbox(slide, Inches(7.1), Inches(2.4), Inches(5.3), Inches(2.5), semantic_lines, 13)

# Analogy
add_shape_box(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.1), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.4), "Real-Life Analogy: Spell Check vs Fact Check", 16, ACCENT_ORANGE, True)
spell_lines = [
    ("Schema Validation = SPELL CHECK: Catches typos, missing words, grammatical errors. But it won't catch a lie!", False, ACCENT_BLUE),
    ("Semantic Validation = FACT CHECK: Catches logical errors. 'The sun orbits the Earth' has perfect spelling but is WRONG.", False, ACCENT_GREEN),
    ("", False, None),
    ("You need BOTH: spell check (schema) catches structural issues, fact check (semantic) catches logic issues.", True, DARK_GRAY),
    ("Example: Schema says 'score must be a number' (OK). Semantic says 'score must be 0-10' (catches score=99).", False, ACCENT_ORANGE),
]
add_multiline_textbox(slide, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.5), spell_lines, 13)


# --- S4 Slide: Retry Loop with Error Feedback ---
slide = content_slide("The Retry Loop -- Self-Correcting Pipeline", "S4")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "When validation fails, DON'T just retry blindly. Feed the error BACK to Claude so it can fix it:", 17, DARK_GRAY, True)

# Flow diagram
add_rounded_box_with_text(slide, Inches(1.0), Inches(1.9), Inches(2.5), Inches(0.7), ACCENT_BLUE, "1. Send Prompt\n+ Tool Schema", 13, WHITE)
add_arrow(slide, Inches(3.5), Inches(2.1), Inches(0.5), Inches(0.3), MED_GRAY)

add_rounded_box_with_text(slide, Inches(4.1), Inches(1.9), Inches(2.5), Inches(0.7), ACCENT_GREEN, "2. Claude Returns\nStructured JSON", 13, WHITE)
add_arrow(slide, Inches(6.6), Inches(2.1), Inches(0.5), Inches(0.3), MED_GRAY)

add_rounded_box_with_text(slide, Inches(7.2), Inches(1.9), Inches(2.5), Inches(0.7), ACCENT_ORANGE, "3. Validate\n(Schema+Semantic)", 13, WHITE)
add_arrow(slide, Inches(9.7), Inches(2.1), Inches(0.5), Inches(0.3), MED_GRAY)

add_rounded_box_with_text(slide, Inches(10.3), Inches(1.9), Inches(2.5), Inches(0.7), ACCENT_PURPLE, "4. Pass?\nUse the Data!", 13, WHITE)

# Fail branch (loops back)
add_down_arrow(slide, Inches(8.25), Inches(2.6), Inches(0.3), Inches(0.4), ACCENT_RED)
add_rounded_box_with_text(slide, Inches(6.7), Inches(3.2), Inches(3.5), Inches(0.7), ACCENT_RED, "FAIL: Send Error Back\nto Claude as tool_result", 12, WHITE)
add_textbox(slide, Inches(6.7), Inches(4.0), Inches(3.5), Inches(0.4), "Claude reads the error and fixes it!", 11, ACCENT_RED, True, PP_ALIGN.CENTER)

# Code example
add_shape_box(slide, Inches(0.6), Inches(4.5), Inches(6.0), Inches(2.8), CODE_BG)
retry_code = [
    ("# Retry loop with error feedback", False, MED_GRAY),
    ("", False, None),
    ("for attempt in range(MAX_RETRIES):", False, WHITE),
    ("    response = call_claude(messages)", False, WHITE),
    ("    result = parse_tool_output(response)", False, WHITE),
    ("", False, None),
    ("    errors = validate(result)", False, ACCENT_ORANGE),
    ("    if not errors:", False, WHITE),
    ("        return result  # Valid!", False, ACCENT_GREEN),
    ("", False, None),
    ("    # Feed error BACK to Claude", False, MED_GRAY),
    ("    messages.append({", False, WHITE),
    ('      "role": "user",', False, WHITE),
    ('      "content": [{', False, WHITE),
    ('        "type": "tool_result",', False, ACCENT_RED),
    ('        "content": f"Error: {errors}.',False, ACCENT_RED),
    ('         Please fix and retry."', False, ACCENT_RED),
    ("      }]", False, WHITE),
    ("    })", False, WHITE),
]
add_multiline_textbox(slide, Inches(0.8), Inches(4.6), Inches(5.6), Inches(2.6), retry_code, 11, WHITE, "Courier New")

# Key points
add_textbox(slide, Inches(7.0), Inches(4.5), Inches(6.0), Inches(0.5), "Key Retry Principles:", 18, DARK_GRAY, True)

retry_points = [
    ("Feed errors back", "Don't just retry silently.\nSend the validation error\nso Claude knows what's wrong.", ACCENT_RED),
    ("Set MAX_RETRIES", "3 retries is typical. If still\nfailing, the source data may\nbe genuinely missing.", ACCENT_ORANGE),
    ("Limits of retry", "Retry can't fix MISSING data.\nIf the source text doesn't\ncontain a name, no retry\nwill produce one.", ACCENT_PURPLE),
]
for i, (title, desc, clr) in enumerate(retry_points):
    y = Inches(5.0 + i * 0.85)
    add_rounded_box_with_text(slide, Inches(7.0), y, Inches(2.3), Inches(0.35), clr, title, 11, WHITE)
    add_textbox(slide, Inches(9.5), y, Inches(3.5), Inches(0.8), desc, 11, DARK_GRAY)


# --- S4 Slide: Validation Analogy ---
slide = content_slide("Real-Life Analogy: Airport Security Checkpoints", "S4")

add_shape_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(6.0), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(1.3), Inches(11), Inches(0.5), "Validation is Like Airport Security -- Multiple Layers:", 22, ACCENT_BLUE, True)

airport_lines = [
    ("", False, None),
    ("Checkpoint 1 -- SCHEMA VALIDATION (= ID + Boarding Pass Check):", True, ACCENT_BLUE),
    ('  Do you have the right documents? Is your name on the boarding pass? Is the flight date correct?', False, MED_GRAY),
    ('  This catches: expired passport, wrong name, missing boarding pass.', False, ACCENT_BLUE),
    ('  This CANNOT catch: whether you\'re actually a threat!', False, ACCENT_RED),
    ("", False, None),
    ("Checkpoint 2 -- SEMANTIC VALIDATION (= X-Ray + Metal Detector):", True, ACCENT_GREEN),
    ('  Is the CONTENT of your bag safe? Does anything violate the rules?', False, MED_GRAY),
    ('  This catches: prohibited items, liquids over 100ml, sharp objects.', False, ACCENT_GREEN),
    ('  These are LOGIC checks your ID can\'t reveal!', False, ACCENT_ORANGE),
    ("", False, None),
    ("Retry Loop (= 'Please remove your belt and try again'):", True, ACCENT_ORANGE),
    ('  When the metal detector beeps, you don\'t get arrested -- you get FEEDBACK and try again.', False, MED_GRAY),
    ('  Same with Claude: validation fails, error is fed back, Claude self-corrects.', False, ACCENT_ORANGE),
    ("", False, None),
    ("MAX_RETRIES = 3 (= After 3 beeps, security pulls you aside for manual review!)", True, ACCENT_RED),
    ("If Claude still fails after retries, the source data may be genuinely problematic.", False, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(1.7), Inches(11), Inches(5.0), airport_lines, 14)


# ═══════════════════════════════════════════════════════════════════
# S5: BATCH PROCESSING WITH CLAUDE API
# ═══════════════════════════════════════════════════════════════════
section_header_slide("S5: Batch Processing with\nClaude API", "The Message Batches API for large offline workloads at ~50% cost savings", ACCENT_PURPLE, "S5")


# --- S5 Slide: What is the Message Batches API? ---
slide = content_slide("The Message Batches API -- Bulk Processing at Scale", "S5")

add_shape_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.5), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
batch_intro = [
    ("Message Batches API = Send thousands of requests at once for asynchronous processing.", True, DARK_GRAY),
    ("~50% cheaper than real-time API calls. Processes within ~24 hours. Perfect for offline workloads.", False, DARK_GRAY),
    ("Think of it as a postal service: you mail a box of letters and pick up the replies the next day.", False, MED_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.3), batch_intro, 16)

# When to use vs not
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5), "When to Use Batch vs Real-Time:", 20, DARK_GRAY, True)

# Batch
add_shape_box(slide, Inches(0.6), Inches(3.6), Inches(5.8), Inches(2.5), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(0.8), Inches(3.7), Inches(5.3), Inches(0.4), "USE Batch API When:", 18, ACCENT_GREEN, True)
use_batch = [
    ("- Classifying 10,000+ support tickets", False, MED_GRAY),
    ("- Analyzing a dataset overnight", False, MED_GRAY),
    ("- Bulk content moderation", False, MED_GRAY),
    ("- Processing log files in bulk", False, MED_GRAY),
    ("- Any workload that can wait hours", False, MED_GRAY),
    ("", False, None),
    ("Key: Non-real-time, high volume,", True, ACCENT_GREEN),
    ("cost-sensitive workloads", True, ACCENT_GREEN),
]
add_multiline_textbox(slide, Inches(0.8), Inches(4.1), Inches(5.3), Inches(2.0), use_batch, 13)

# Real-time
add_shape_box(slide, Inches(6.9), Inches(3.6), Inches(5.8), Inches(2.5), RGBColor(0xFA, 0xEB, 0xEB), ACCENT_RED, Pt(2))
add_textbox(slide, Inches(7.1), Inches(3.7), Inches(5.3), Inches(0.4), "DON'T Use Batch When:", 18, ACCENT_RED, True)
no_batch = [
    ("- User is waiting for a response", False, MED_GRAY),
    ("- Interactive chat conversations", False, MED_GRAY),
    ("- Multi-turn tool use within a batch", False, MED_GRAY),
    ("- Real-time code review in IDE", False, MED_GRAY),
    ("- Latency-sensitive applications", False, MED_GRAY),
    ("", False, None),
    ("Key: Interactive, real-time, or", True, ACCENT_RED),
    ("multi-turn tool use scenarios", True, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(7.1), Inches(4.1), Inches(5.3), Inches(2.0), no_batch, 13)

# custom_id
add_shape_box(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.9), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.8),
    "KEY: Each batch request needs a custom_id (e.g., 'ticket-001'). When results come back, match custom_id to your input.\nHandle failures by resubmitting only the failed items, not the entire batch!", 14, ACCENT_ORANGE, True)


# --- S5 Slide: Batch Processing Code Example ---
slide = content_slide("Batch API -- Code Example & Workflow", "S5")

# Code example
add_shape_box(slide, Inches(0.6), Inches(1.2), Inches(6.5), Inches(4.5), CODE_BG)
batch_code = [
    ("# Message Batches API workflow", False, MED_GRAY),
    ("", False, None),
    ("import anthropic", False, WHITE),
    ("client = anthropic.Anthropic()", False, WHITE),
    ("", False, None),
    ("# 1. Prepare batch requests", False, MED_GRAY),
    ("requests = [", False, WHITE),
    ("  {", False, WHITE),
    ('    "custom_id": "ticket-001",', True, ACCENT_GREEN),
    ('    "params": {', False, WHITE),
    ('      "model": "claude-sonnet-4-6",', False, ACCENT_BLUE),
    ('      "max_tokens": 1024,', False, WHITE),
    ('      "messages": [{', False, WHITE),
    ('        "role": "user",', False, WHITE),
    ('        "content": "Classify: ..."', False, WHITE),
    ("      }]", False, WHITE),
    ("    }", False, WHITE),
    ("  },", False, WHITE),
    ("  # ... 10,000 more requests", False, MED_GRAY),
    ("]", False, WHITE),
    ("", False, None),
    ("# 2. Submit the batch", False, MED_GRAY),
    ("batch = client.messages.batches", False, WHITE),
    ("  .create(requests=requests)", False, WHITE),
    ("", False, None),
    ("# 3. Poll for results", False, MED_GRAY),
    ("# Results ready within ~24 hours", False, MED_GRAY),
]
add_multiline_textbox(slide, Inches(0.8), Inches(1.3), Inches(6.1), Inches(4.3), batch_code, 11, WHITE, "Courier New")

# Workflow steps
add_textbox(slide, Inches(7.5), Inches(1.2), Inches(5.5), Inches(0.5), "Batch Processing Workflow:", 20, DARK_GRAY, True)

batch_steps = [
    ("1. Prepare Requests", "Bundle all items with unique\ncustom_id per request.\nKeep track of IDs locally.", ACCENT_BLUE),
    ("2. Submit Batch", "One API call sends all requests.\n~50% cheaper than real-time.\nAsync -- returns immediately.", ACCENT_GREEN),
    ("3. Poll for Status", "Check batch status periodically.\nResults ready in up to ~24h.\nPartial results available early.", ACCENT_ORANGE),
    ("4. Handle Failures", "Some items may fail. Resubmit\nonly failed items by custom_id.\nDon't re-run the whole batch!", ACCENT_RED),
]
for i, (title, desc, clr) in enumerate(batch_steps):
    y = Inches(1.8 + i * 1.3)
    add_rounded_box_with_text(slide, Inches(7.5), y, Inches(2.5), Inches(0.4), clr, title, 12, WHITE)
    add_textbox(slide, Inches(7.5), y + Inches(0.5), Inches(5.5), Inches(0.8), desc, 12, DARK_GRAY)

# Analogy
add_shape_box(slide, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.3), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(6.1), Inches(11), Inches(0.4), "Real-Life Analogy: Laundromat vs Hand Washing", 16, ACCENT_ORANGE, True)
laundry_lines = [
    ("Real-time API = Washing each shirt by hand. Fast per item, but exhausting at scale.", False, ACCENT_RED),
    ("Batch API = Loading the washing machine with 50 shirts. Takes longer, but way cheaper and easier!", False, ACCENT_GREEN),
    ("You wouldn't hand-wash 10,000 shirts. Same logic: use Batch API for large-scale processing.", True, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.7), laundry_lines, 13)


# ═══════════════════════════════════════════════════════════════════
# S6: MULTI-INSTANCE & MULTI-PASS ARCHITECTURES
# ═══════════════════════════════════════════════════════════════════
section_header_slide("S6: Multi-Instance & Multi-Pass\nArchitectures", "Why self-review is biased, and how independent instances + multi-pass improve quality", ACCENT_TEAL, "S6")


# --- S6 Slide: Why Self-Review is Biased ---
slide = content_slide("Why Self-Review is Biased -- The Anchoring Problem", "S6")

add_shape_box(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.5), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
bias_intro = [
    ("When Claude reviews its OWN output, it's biased -- it remembers what it wrote and why.", True, DARK_GRAY),
    ("This is called 'anchoring bias'. The model is anchored to its original reasoning and tends to confirm it.", False, DARK_GRAY),
    ("Solution: Use a SEPARATE, independent Claude instance that has never seen the original generation.", False, MED_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.3), bias_intro, 16)

# Self-review vs Independent
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.5), "Self-Review vs Independent Instance Review:", 20, DARK_GRAY, True)

# Self-review (BAD)
add_shape_box(slide, Inches(0.6), Inches(3.6), Inches(5.8), Inches(2.5), RGBColor(0xFA, 0xEB, 0xEB), ACCENT_RED, Pt(2))
add_textbox(slide, Inches(0.8), Inches(3.7), Inches(5.3), Inches(0.4), "SELF-REVIEW (Biased)", 18, ACCENT_RED, True)
self_review = [
    ("Same Claude instance writes AND reviews.", False, DARK_GRAY),
    ("", False, None),
    ("Problems:", True, ACCENT_RED),
    ("- Remembers original reasoning", False, MED_GRAY),
    ("- Tends to confirm its own output", False, MED_GRAY),
    ("- Misses errors it made initially", False, MED_GRAY),
    ("- Like grading your own exam!", False, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(0.8), Inches(4.1), Inches(5.3), Inches(2.0), self_review, 14)

# Independent (GOOD)
add_shape_box(slide, Inches(6.9), Inches(3.6), Inches(5.8), Inches(2.5), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(7.1), Inches(3.7), Inches(5.3), Inches(0.4), "INDEPENDENT INSTANCE (Unbiased)", 18, ACCENT_GREEN, True)
indep_review = [
    ("Fresh Claude instance reviews output.", False, DARK_GRAY),
    ("", False, None),
    ("Benefits:", True, ACCENT_GREEN),
    ("- No memory of original reasoning", False, MED_GRAY),
    ("- Evaluates output on its own merit", False, MED_GRAY),
    ("- Catches blind spots and errors", False, MED_GRAY),
    ("- Like peer review by a colleague!", False, ACCENT_GREEN),
]
add_multiline_textbox(slide, Inches(7.1), Inches(4.1), Inches(5.3), Inches(2.0), indep_review, 14)

# Analogy
add_shape_box(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.9), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.8),
    "Real-Life: An author ALWAYS has a separate editor review their book. You can't proofread your own work well --\nyour brain auto-corrects because it KNOWS what you meant to write. Same problem with AI self-review!", 14, ACCENT_ORANGE, True)


# --- S6 Slide: Multi-Pass Architecture ---
slide = content_slide("Multi-Pass Architecture -- Per-Item + Cross-Item", "S6")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "Multi-pass = Multiple rounds of analysis, each building on the last. Two key patterns:", 17, DARK_GRAY, True)

# Pass 1
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(2.8), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(0.4), "Pass 1: Per-Item Analysis", 20, ACCENT_BLUE, True)
pass1_lines = [
    ("Each item reviewed INDEPENDENTLY.", True, DARK_GRAY),
    ("", False, None),
    ("How it works:", True, ACCENT_BLUE),
    ("- Item 1 -> Claude -> Result 1", False, MED_GRAY),
    ("- Item 2 -> Claude -> Result 2", False, MED_GRAY),
    ("- Item 3 -> Claude -> Result 3", False, MED_GRAY),
    ("- (Can run in parallel!)", False, ACCENT_GREEN),
    ("", False, None),
    ("Good at: detailed per-item analysis", True, DARK_GRAY),
    ("Missing: cross-item patterns!", True, ACCENT_RED),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.3), Inches(2.2), pass1_lines, 13)

# Arrow between passes
add_arrow(slide, Inches(6.4), Inches(3.2), Inches(0.5), Inches(0.3), ACCENT_ORANGE)

# Pass 2
add_shape_box(slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(2.8), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.4), "Pass 2: Cross-Item Integration", 20, ACCENT_GREEN, True)
pass2_lines = [
    ("ALL results reviewed TOGETHER.", True, DARK_GRAY),
    ("", False, None),
    ("How it works:", True, ACCENT_GREEN),
    ("- Collect all Pass 1 results", False, MED_GRAY),
    ("- Feed ALL to a new Claude call", False, MED_GRAY),
    ("- Find patterns across items", False, MED_GRAY),
    ("- Resolve contradictions", False, MED_GRAY),
    ("", False, None),
    ("Good at: finding trends, inconsistencies", True, DARK_GRAY),
    ("and systemic issues!", True, ACCENT_GREEN),
]
add_multiline_textbox(slide, Inches(7.1), Inches(2.4), Inches(5.3), Inches(2.2), pass2_lines, 13)

# Pipeline diagram
add_shape_box(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.3), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(5.1), Inches(11), Inches(0.4), "Real-Life Analogy: How Doctors Diagnose Complex Cases", 16, ACCENT_ORANGE, True)
doctor_lines = [
    ("Pass 1 (Per-Item) = Each SPECIALIST examines the patient independently:", False, ACCENT_BLUE),
    ("  Cardiologist checks heart, neurologist checks brain, radiologist reads scans -- each in isolation.", False, MED_GRAY),
    ("", False, None),
    ("Pass 2 (Cross-Item) = The TEAM MEETING where all specialists share findings:", False, ACCENT_GREEN),
    ("  'The heart issue + the brain scan + the blood work together suggest X'. Patterns emerge from integration!", False, MED_GRAY),
    ("", False, None),
    ("Single-pass = asking ONE doctor to do everything. Multi-pass = a specialist team. Quality goes up!", True, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.7), doctor_lines, 13)


# --- S6 Slide: Parallelization & Chaining ---
slide = content_slide("Parallelization vs Chaining Workflows", "S6")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
    "Multi-instance uses two workflow patterns from the Claude API:", 17, DARK_GRAY, True)

# Parallelization
add_shape_box(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(3.0), RGBColor(0xEA, 0xF2, 0xF8), ACCENT_BLUE, Pt(2))
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.3), Inches(0.4), "Parallelization (Fan-Out)", 20, ACCENT_BLUE, True)
parallel_lines = [
    ("Multiple Claude instances run", True, DARK_GRAY),
    ("AT THE SAME TIME on different items.", False, DARK_GRAY),
    ("", False, None),
    ("        Input Data", False, ACCENT_BLUE),
    ("       /    |    \\", False, MED_GRAY),
    ("   Claude  Claude  Claude", False, ACCENT_BLUE),
    ("       \\    |    /", False, MED_GRAY),
    ("      Collect Results", False, ACCENT_GREEN),
    ("", False, None),
    ("Use for: independent tasks,", True, DARK_GRAY),
    ("per-item analysis, throughput", False, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.3), Inches(2.5), parallel_lines, 13, DARK_GRAY, "Courier New")

# Chaining
add_shape_box(slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(3.0), RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN, Pt(2))
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.4), "Chaining (Sequential Passes)", 20, ACCENT_GREEN, True)
chaining_lines = [
    ("Output of one Claude call becomes", True, DARK_GRAY),
    ("INPUT to the next Claude call.", False, DARK_GRAY),
    ("", False, None),
    ("  Claude Pass 1 (Draft)", False, ACCENT_BLUE),
    ("         |", False, MED_GRAY),
    ("         v", False, MED_GRAY),
    ("  Claude Pass 2 (Critique)", False, ACCENT_GREEN),
    ("         |", False, MED_GRAY),
    ("         v", False, MED_GRAY),
    ("  Claude Pass 3 (Refine)", False, ACCENT_ORANGE),
    ("", False, None),
    ("Use for: draft-critique-refine,", True, DARK_GRAY),
    ("multi-pass review, quality gates", False, DARK_GRAY),
]
add_multiline_textbox(slide, Inches(7.1), Inches(2.4), Inches(5.3), Inches(2.5), chaining_lines, 13, DARK_GRAY, "Courier New")

# When to combine
add_shape_box(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.1), RGBColor(0xFD, 0xF2, 0xE9), ACCENT_ORANGE, Pt(2))
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(11), Inches(0.4), "Combining Both: The Power Pattern", 18, ACCENT_ORANGE, True)
combine_lines = [
    ("Step 1 (Parallel): Fan-out 100 items to 100 Claude instances. Each analyzes one item independently.", False, ACCENT_BLUE),
    ("Step 2 (Chain): Feed ALL 100 results to a new Claude instance for cross-item synthesis and report.", False, ACCENT_GREEN),
    ("", False, None),
    ("This is the multi-pass architecture: Parallelization for Pass 1 (per-item), Chaining for Pass 2 (cross-item).", True, DARK_GRAY),
    ("Analogy: 100 detectives investigate 100 leads (parallel), then ONE lead detective writes the final case report (chain).", False, ACCENT_ORANGE),
]
add_multiline_textbox(slide, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.5), combine_lines, 13)


# ═══════════════════════════════════════════════════════════════════
# EXAM PREP: Key Questions
# ═══════════════════════════════════════════════════════════════════
slide = content_slide("Exam Prep -- Key Questions to Answer")

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5), "Can you answer these from your Module 4 study?", 20, DARK_GRAY, True)

questions = [
    ("Q1", "What makes a prompt's classification criteria explicit (categorical rules + severity tiers),\nand how does that cut false positives?", ACCENT_RED),
    ("Q2", "When do few-shot examples help, how many do you use, and how do you get the model\nto generalise beyond them?", ACCENT_BLUE),
    ("Q3", "How do tool_use + a JSON schema force structured output, what does tool_choice: any do,\nand what can a schema NOT catch?", ACCENT_GREEN),
    ("Q4", "How do schema (syntax) and semantic (logic) validation differ, and how does an\nerror-feedback retry loop self-correct?", ACCENT_ORANGE),
    ("Q5", "When is the Message Batches API the right choice (and when not)? What is custom_id for?\nHow do you handle failures?", ACCENT_PURPLE),
    ("Q6", "Why use an independent instance instead of self-review? How does multi-pass (per-item\nthen cross-item) improve quality?", ACCENT_TEAL),
]

for i, (q, text, clr) in enumerate(questions):
    y = Inches(1.7 + i * 0.95)
    add_rounded_box_with_text(slide, Inches(0.6), y, Inches(0.7), Inches(0.5), clr, q, 14, WHITE)
    add_textbox(slide, Inches(1.5), y, Inches(11.3), Inches(0.9), text, 14, DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════
# FINAL SLIDE: Cheat Sheet / Exam Prep
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8), "Module 4 -- Exam Cheat Sheet", 36, WHITE, True, PP_ALIGN.CENTER)

cheat_items = [
    ("S1", "Explicit Criteria", "Use categorical rules + severity tiers + concrete examples. Iterate to reduce false positives. Vague prompts = inconsistent results.", ACCENT_RED),
    ("S2", "Few-Shot Prompting", "2-4 examples that cover AMBIGUOUS cases. Lock in format + tone. Claude generalizes the pattern to unseen inputs.", ACCENT_BLUE),
    ("S3", "tool_use + JSON Schema", "tool_choice: any forces structured output. enum fields restrict values. nullable for missing data. Schema catches syntax NOT logic.", ACCENT_GREEN),
    ("S4", "Validation & Retry", "Schema = syntax check, Semantic = logic check. Feed errors BACK to Claude in retry loop. MAX_RETRIES ~3. Can't fix missing data.", ACCENT_ORANGE),
    ("S5", "Batch Processing", "Message Batches API: ~50% cheaper, async ~24h. Use custom_id per request. NOT for interactive/multi-turn. Resubmit failures only.", ACCENT_PURPLE),
    ("S6", "Multi-Instance/Pass", "Self-review is biased (anchoring). Use independent instance. Pass 1 = per-item (parallel), Pass 2 = cross-item (chain).", ACCENT_TEAL),
]

for i, (sec, title, key, clr) in enumerate(cheat_items):
    y = Inches(1.2 + i * 1.0)
    add_rounded_box_with_text(slide, Inches(0.5), y, Inches(0.7), Inches(0.55), clr, sec, 14, WHITE)
    add_textbox(slide, Inches(1.4), y, Inches(2.8), Inches(0.35), title, 18, WHITE, True)
    add_textbox(slide, Inches(1.4), y + Inches(0.35), Inches(11.3), Inches(0.55), key, 13, RGBColor(0xBB, 0xBB, 0xBB))

# Key reminder
add_shape_box(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3), RGBColor(0x2A, 0x2A, 0x4E))
add_textbox(slide, Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.4),
    "Remember: Explicit criteria = consistency | Few-shot = format lock | tool_use = guaranteed JSON | Retry = self-correct | Batch = scale | Multi-pass = quality", 12, YELLOW, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
output_path = "/Users/rohit.sakate/Documents/mediware/cca-f/week-4/CCA-F_Module4_Prompt_Engineering_Structured_Output.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
